from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUTS = ROOT / "outputs"
FIGURES = OUTPUTS / "figures"
ASSETS = ROOT / "assets"
DECK_ASSETS = FIGURES / "deck-assets"

PPTX_OUT = OUTPUTS / "2026-06-11_presentacion-sueno-memoria_v2.pptx"
SCRIPT_OUT = OUTPUTS / "2026-06-11_guion-presentacion-sueno-memoria_v2.md"

W = Inches(13.333)
H = Inches(7.5)

INK = RGBColor(27, 38, 59)
MUTED = RGBColor(82, 95, 111)
PAPER = RGBColor(248, 246, 241)
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(76, 111, 148)
ORANGE = RGBColor(206, 113, 49)
GREEN = RGBColor(80, 129, 111)
GRAY = RGBColor(226, 232, 240)
PALE_BLUE = RGBColor(231, 238, 247)
PALE_ORANGE = RGBColor(252, 237, 220)


def add_bg(slide, color=PAPER):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(
    slide,
    text,
    x,
    y,
    w,
    h,
    size=24,
    bold=False,
    color=INK,
    align=PP_ALIGN.LEFT,
    font="Aptos",
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, bullets, x, y, w, h, size=18, color=INK, gap=8):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
    return box


def add_kicker(slide, label, idx):
    marker = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.39), Inches(0.28), Inches(0.55), Inches(0.36))
    marker.fill.solid()
    marker.fill.fore_color.rgb = ORANGE
    marker.line.fill.background()
    add_text(slide, f"{idx:02d}", Inches(0.45), Inches(0.32), Inches(0.42), Inches(0.28), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    add_text(slide, label.upper(), Inches(1.05), Inches(0.3), Inches(3.8), Inches(0.3), size=9, bold=True, color=MUTED, valign=MSO_ANCHOR.MIDDLE)


def add_title(slide, idx, kicker, title, claim=None):
    add_kicker(slide, kicker, idx)
    add_text(slide, title, Inches(0.55), Inches(0.85), Inches(7.7), Inches(0.55), size=25, bold=True, color=INK)
    if claim:
        add_text(slide, claim, Inches(0.58), Inches(1.45), Inches(8.2), Inches(0.55), size=13, color=MUTED)


def add_rule(slide, y=Inches(6.88)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), y, Inches(12.2), Pt(0.7))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(210, 216, 224)
    line.line.fill.background()


def add_source(slide, text):
    add_rule(slide)
    add_text(slide, text, Inches(0.58), Inches(6.96), Inches(11.8), Inches(0.22), size=7.5, color=MUTED)


def add_chip(slide, text, x, y, color, width=1.65):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(width), Inches(0.38))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text(slide, text, x + Inches(0.1), y + Inches(0.06), Inches(width - 0.2), Inches(0.2), size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    return shape


def add_card(slide, x, y, w, h, title, body, accent=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(220, 226, 234)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, title, x + Inches(0.22), y + Inches(0.18), w - Inches(0.38), Inches(0.3), size=14, bold=True, color=INK)
    add_text(slide, body, x + Inches(0.22), y + Inches(0.58), w - Inches(0.42), h - Inches(0.68), size=11.5, color=MUTED)
    return shape


def add_picture_fit(slide, path: Path, x, y, w, h):
    from PIL import Image

    path = prepare_deck_image(path)
    with Image.open(path) as img:
        iw, ih = img.size
    ratio = iw / ih
    box_ratio = w / h
    if ratio > box_ratio:
        pic_w = w
        pic_h = w / ratio
        pic_x = x
        pic_y = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * ratio
        pic_x = x + (w - pic_w) / 2
        pic_y = y
    return slide.shapes.add_picture(str(path), pic_x, pic_y, width=pic_w, height=pic_h)


def prepare_deck_image(path: Path, max_edge=2200) -> Path:
    from PIL import Image

    DECK_ASSETS.mkdir(parents=True, exist_ok=True)
    suffix = ".jpg" if path.parent == ASSETS else ".png"
    out = DECK_ASSETS / f"{path.stem}{suffix}"
    with Image.open(path) as img:
        img = img.convert("RGBA" if suffix == ".png" else "RGB")
        img.thumbnail((max_edge, max_edge), Image.Resampling.LANCZOS)
        if suffix == ".jpg":
            img.save(out, quality=86, optimize=True)
        else:
            img.save(out, optimize=True)
    return out


def add_table(slide, rows, x, y, w, h, col_widths=None, font_size=10):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), x, y, w, h)
    table = table_shape.table
    if col_widths:
        for idx, col_w in enumerate(col_widths):
            table.columns[idx].width = col_w
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(value)
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            fill = cell.fill
            fill.solid()
            fill.fore_color.rgb = INK if r == 0 else (RGBColor(245, 247, 250) if r % 2 else WHITE)
            for p in cell.text_frame.paragraphs:
                p.font.name = "Aptos"
                p.font.size = Pt(font_size)
                p.font.bold = r == 0
                p.font.color.rgb = WHITE if r == 0 else INK
    return table_shape


def new_slide(prs, idx, kicker, title, claim=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, idx, kicker, title, claim)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    slides_script = []

    # 0
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, INK)
    add_picture_fit(s, ASSETS / "pre-siesta-listo.png", Inches(7.0), Inches(0.0), Inches(6.33), H)
    add_text(s, "Sueño y consolidacion de memoria declarativa", Inches(0.65), Inches(0.78), Inches(6.1), Inches(1.35), size=29, bold=True, color=WHITE)
    add_text(s, "Informe tipo paper en formato PPTX | Neurociencias del Desarrollo Productivo", Inches(0.68), Inches(2.35), Inches(5.55), Inches(0.55), size=14, color=RGBColor(222, 229, 238))
    add_chip(s, "v2 - 2026-06-11", Inches(0.68), Inches(3.05), ORANGE, width=1.75)
    add_text(s, "Equipo: completar integrantes\nFecha de exposicion: completar", Inches(0.68), Inches(5.85), Inches(5.7), Inches(0.7), size=12, color=RGBColor(222, 229, 238))
    slides_script.append(("0", "Abrir como si fuera un paper oral: pregunta, hipotesis, metodo, resultados, discusion y conclusion. Aclarar que la version v2 prioriza integracion con clases teoricas."))

    # 1
    s = new_slide(prs, 1, "abstract", "Resumen del estudio", "Hipotesis central: el sueño mejora la consolidacion de memoria declarativa.")
    add_card(s, Inches(0.75), Inches(1.9), Inches(3.7), Inches(3.0), "Pregunta", "¿Dormir despues del aprendizaje mejora la retencion de palabras y definiciones?", BLUE)
    add_card(s, Inches(4.8), Inches(1.9), Inches(3.7), Inches(3.0), "Resultado", "El sujeto en condicion sueño tuvo mayor recuerdo formal de palabras; las definiciones fueron altas en casi todos.", GREEN)
    add_card(s, Inches(8.85), Inches(1.9), Inches(3.7), Inches(3.0), "Conclusion", "La hipotesis queda apoyada descriptivamente, pero no confirmada estadisticamente por n pequeño y fuentes separadas.", ORANGE)
    add_text(s, "Mensaje de paper: evidencia compatible, conclusion cauta.", Inches(1.0), Inches(5.45), Inches(11.3), Inches(0.45), size=18, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_source(s, "Fuentes: consigna oficial; outputs/2026-06-11_resultados-memoria_v1.md.")
    slides_script.append(("1", "Dar el resumen en 45 segundos: hipotesis, patron observado y conclusion cauta."))

    # 2
    s = new_slide(prs, 2, "introduccion", "Introduccion teorica", "El sueño favorece memoria cuando su arquitectura coordina fases, ritmos y neuromodulacion.")
    add_table(
        s,
        [
            ["Concepto de clase", "Rol en el TP"],
            ["NREM/SWS", "contexto fisiologico para consolidacion declarativa"],
            ["Husos 12-15 Hz", "oscilaciones N2/NREM asociadas a memoria"],
            ["Ondas lentas + ripples", "marco temporal para transferencia hipocampo-corteza"],
            ["REM", "relevante para otros tipos de memoria; no eje principal aqui"],
        ],
        Inches(0.8),
        Inches(1.95),
        Inches(7.15),
        Inches(2.75),
        [Inches(2.4), Inches(4.75)],
        font_size=11.5,
    )
    add_bullets(
        s,
        [
            "La memoria declarativa depende de circuitos hipocampales.",
            "Dormir no solo evita interferencia: puede reactivar y reorganizar memorias.",
            "La calidad de ritmos importa mas que decir simplemente 'durmio'.",
        ],
        Inches(8.35),
        Inches(2.0),
        Inches(4.0),
        Inches(2.6),
        size=14,
    )
    add_source(s, "Fuentes: clase 3 sueño y memoria; clase 8 farmacos/husos; clase 9 EEG.")
    slides_script.append(("2", "Enmarcar el trabajo con clases teoricas: arquitectura del sueño, SWS, husos y reactivacion hipocampo-cortical."))

    # 3
    s = new_slide(prs, 3, "marco", "Modelo mecanistico", "Consolidacion activa: ondas lentas, husos y ripples coordinan la reactivacion.")
    add_card(s, Inches(0.75), Inches(2.0), Inches(3.65), Inches(2.65), "Hipocampo", "reactiva trazas recientes y secuencias de memoria", BLUE)
    add_card(s, Inches(4.85), Inches(2.0), Inches(3.65), Inches(2.65), "Oscilaciones", "ondas lentas corticales + husos talamocorticales + ripples", GREEN)
    add_card(s, Inches(8.95), Inches(2.0), Inches(3.65), Inches(2.65), "Neocorteza", "integra y estabiliza representaciones declarativas", ORANGE)
    add_text(s, "Prediccion: si el intervalo incluye sueño fisiologicamente organizado, el recuerdo posterior deberia mejorar.", Inches(1.1), Inches(5.35), Inches(11.0), Inches(0.5), size=17, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_source(s, "Fuentes: clase 3; clase 8 sobre acoplamiento ondas lentas-husos-ripples.")
    slides_script.append(("3", "Explicar el mecanismo: no es magia ni descanso general; es reactivacion organizada por oscilaciones."))

    # 4
    s = new_slide(prs, 4, "hipotesis", "Hipotesis y criterio de decision", "H1: dormir mejora la consolidacion de la tarea de palabras.")
    add_table(
        s,
        [
            ["Prediccion", "Resultado esperado si H1 se apoya"],
            ["Conductual", "mayor recuerdo de palabra objetivo y/o definicion tras sueño"],
            ["Fisiologico", "registro con fases NREM y ritmos compatibles con sueño"],
            ["Mecanistico", "presencia esperable de SWS/N2, husos o sigma como marco teorico"],
            ["Decision", "aceptacion descriptiva si patron favorece sueño; no inferencia fuerte con n=1"],
        ],
        Inches(0.75),
        Inches(1.95),
        Inches(11.8),
        Inches(3.25),
        [Inches(2.55), Inches(9.25)],
        font_size=11.5,
    )
    add_text(s, "Criterio final: apoyada parcialmente/descriptivamente, rechazada si sueño no mejora o si no hay firma fisiologica compatible.", Inches(0.95), Inches(5.65), Inches(11.4), Inches(0.45), size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_source(s, "Fuentes: consigna oficial; clase 3 sueño y memoria.")
    slides_script.append(("4", "Dejar claro que aceptar/rechazar no es estadistico en sentido fuerte: es una conclusion descriptiva del practico."))

    # 5
    s = new_slide(prs, 5, "metodos", "Metodos: diseño y procedencia", "El analisis integra conducta y fisiologia, pero no mezcla sujetos ni instrumentos.")
    add_table(
        s,
        [
            ["Fuente", "Uso", "Restriccion"],
            ["Planilla memoria", "resultado conductual por sujeto/condicion", "sueño n=1"],
            ["OpenBCI propio", "evento de no conciliacion", "no se usa para estadificar sueño"],
            ["BrainVision secundario", "hipnograma/espectrograma", "no se correlaciona directo con memoria"],
        ],
        Inches(0.75),
        Inches(1.95),
        Inches(11.8),
        Inches(2.7),
        [Inches(2.4), Inches(4.7), Inches(4.7)],
        font_size=11.5,
    )
    add_text(s, "Regla metodologica: el EEG con sueño explica fisiologia; la planilla explica memoria. La relacion se discute teoricamente, no como correlacion sujeto-a-sujeto.", Inches(1.0), Inches(5.3), Inches(11.2), Inches(0.65), size=14, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_source(s, "Fuentes: README.md, data/README.md, AGENTS.md.")
    slides_script.append(("5", "Esta slide protege el trabajo: paper-like, claro con fuentes y limites de inferencia."))

    # 6
    s = new_slide(prs, 6, "metodos", "Variables y medidas", "Se operacionalizo consolidacion como rendimiento posterior en la tarea de palabras.")
    add_card(s, Inches(0.75), Inches(1.95), Inches(3.65), Inches(2.85), "Variable independiente", "Condicion del intervalo: vigilia vs sueño.", BLUE)
    add_card(s, Inches(4.85), Inches(1.95), Inches(3.65), Inches(2.85), "Variables dependientes", "Palabra objetivo: forma lexical.\nDefinicion: contenido semantico.", GREEN)
    add_card(s, Inches(8.95), Inches(1.95), Inches(3.65), Inches(2.85), "Variables moduladoras", "Cronotipo, estres, cafeina, pantallas, ambiente, adaptacion al dispositivo.", ORANGE)
    add_text(s, "Separar palabra y definicion permite ver si el sueño se asocia mas con forma exacta o con significado.", Inches(1.05), Inches(5.45), Inches(11.0), Inches(0.45), size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_source(s, "Fuente: analysis/analyze_memory.py; clases 6, 7 y 8.")
    slides_script.append(("6", "Explicar las variables como en un paper: independiente, dependientes y moduladores."))

    # 7
    s = new_slide(prs, 7, "metodos", "Scoring de memoria", "Dos puntajes capturan niveles distintos de recuperacion declarativa.")
    add_card(s, Inches(0.9), Inches(2.05), Inches(5.4), Inches(2.65), "Palabra objetivo", "Coincidencia formal exacta tras normalizar mayusculas, acentos y signos. Mide recuperacion lexical/formal.", BLUE)
    add_card(s, Inches(6.95), Inches(2.05), Inches(5.4), Inches(2.65), "Definicion", "Scoring semantico manual binario, item por item. Mide recuperacion del significado.", GREEN)
    add_text(s, "Esto evita tratar como iguales el recuerdo de 'que significaba' y el recuerdo de 'como se llamaba'.", Inches(1.0), Inches(5.35), Inches(11.25), Inches(0.45), size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_source(s, "Fuente: outputs/2026-06-11_resultados-memoria_v1.md.")
    slides_script.append(("7", "Presentar el scoring como metodo, no como resultado."))

    # 8
    s = new_slide(prs, 8, "resultados", "Resultados conductuales", "El patron observado favorece al sujeto en condicion sueño en palabra objetivo.")
    add_picture_fit(s, FIGURES / "2026-06-11_memoria_scores_v1.png", Inches(0.7), Inches(1.78), Inches(12.0), Inches(4.65))
    add_source(s, "Fuente: data/PRACTICO LABORATORIO EQUIPO 1.xlsx; scoring reproducible en analysis/analyze_memory.py.")
    slides_script.append(("8", "Mostrar el grafico y decir: el resultado central para la hipotesis es palabra objetivo 17/20 en sueño vs promedio 5/20 en vigilia."))

    # 9
    s = new_slide(prs, 9, "decision", "Decision sobre la hipotesis", "H1 se apoya descriptivamente, pero no se acepta como evidencia estadistica concluyente.")
    add_table(
        s,
        [
            ["Evidencia", "Lectura"],
            ["Palabra objetivo", "Sueño 17/20; vigilia promedio 5/20"],
            ["Definicion", "Sueño 20/20; vigilia promedio 17.3/20"],
            ["Tamaño muestral", "Sueño n=1; no inferencia estadistica fuerte"],
            ["Conclusion", "Hipotesis apoyada parcialmente/descriptivamente"],
        ],
        Inches(0.8),
        Inches(1.95),
        Inches(7.1),
        Inches(3.15),
        [Inches(2.2), Inches(4.9)],
        font_size=11.5,
    )
    add_card(s, Inches(8.35), Inches(2.0), Inches(4.0), Inches(2.85), "Respuesta tipo paper", "No se rechaza H1. Los datos son compatibles con H1, pero no alcanzan para aceptarla en sentido estadistico fuerte.", ORANGE)
    add_source(s, "Fuente: outputs/2026-06-11_resultados-memoria_v1.md.")
    slides_script.append(("9", "Usar lenguaje de paper: no rechazar, apoyar descriptivamente, no sobreafirmar."))

    # 10
    s = new_slide(prs, 10, "metodos eeg", "Analisis fisiologico secundario", "El EEG con sueño se usa para discutir mecanismos, no para correlacion sujeto-a-sujeto.")
    add_table(
        s,
        [
            ["Registro", "Detalle"],
            ["Canales", "EOG I/D; F3/F4; C3/C4; P3/P4; EMGI/EMGD"],
            ["Muestreo", "250 Hz; 79.78 min"],
            ["Scoring", "159 epocas de 30 s; codigos 0-3 tentativos"],
            ["Marcadores", "luz off 1.05 min; luz on 79.68 min"],
        ],
        Inches(0.75),
        Inches(1.95),
        Inches(7.35),
        Inches(3.15),
        [Inches(2.2), Inches(5.15)],
        font_size=11.5,
    )
    add_picture_fit(s, ASSETS / "Colocando-electrodos.png", Inches(8.45), Inches(2.0), Inches(3.8), Inches(2.65))
    add_source(s, "Fuentes: data/S3practica.vhdr, .vmrk, .eeg; data/S3PRACTICA.txt.")
    slides_script.append(("10", "Aclarar que este registro permite mostrar arquitectura y señales esperables del sueño, pero no calcular correlacion directa con la planilla."))

    # 11
    s = new_slide(prs, 11, "resultados eeg", "Hipnograma tentativo", "El dataset secundario muestra una secuencia compatible con sueño NREM.")
    add_picture_fit(s, FIGURES / "2026-06-11_eeg_hipnograma_v1.png", Inches(0.65), Inches(1.75), Inches(12.1), Inches(4.35))
    add_source(s, "Fuente: data/S3PRACTICA.txt + data/S3practica.vmrk. Scoring tentativo, no clinico.")
    slides_script.append(("11", "Enfatizar el bloque compatible con N3/SWS probable, sin convertirlo en diagnostico clinico."))

    # 12
    s = new_slide(prs, 12, "resultados eeg", "Espectrograma C3/C4", "La señal permite visualizar estructura tiempo-frecuencia durante el intervalo de sueño.")
    add_picture_fit(s, FIGURES / "2026-06-11_eeg_espectrograma_c3_c4_v1.png", Inches(0.65), Inches(1.72), Inches(12.1), Inches(4.45))
    add_source(s, "Fuente: data/S3practica.eeg; filtro 0.3-35 Hz; espectrograma Hann/Welch.")
    slides_script.append(("12", "Usar el espectrograma como resultado fisiologico descriptivo; no como prueba clinica de fases."))

    # 13
    s = new_slide(prs, 13, "husos", "Husos, sigma y memoria", "La correlacion husos-memoria es teoricamente relevante, pero este TP no permite estimarla directamente.")
    add_card(s, Inches(0.75), Inches(1.95), Inches(3.65), Inches(2.9), "Teoria de clase", "Husos de 12-15 Hz se asocian a memoria declarativa, sobre todo acoplados a ondas lentas.", BLUE)
    add_card(s, Inches(4.85), Inches(1.95), Inches(3.65), Inches(2.9), "Dato disponible", "Se calculo banda sigma 12-15 Hz en C3/C4 como descriptor espectral, no detector formal de husos.", GREEN)
    add_card(s, Inches(8.95), Inches(1.95), Inches(3.65), Inches(2.9), "Limite", "No se calcula correlacion husos-memoria porque el EEG con sueño es secundario y no debe mezclarse con la planilla.", ORANGE)
    add_text(s, "Analisis futuro: detector de husos por epoca + rendimiento del mismo sujeto/sesion.", Inches(1.1), Inches(5.45), Inches(11.0), Inches(0.45), size=16, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_source(s, "Fuentes: clase 3; clase 8, seccion sueño, husos y memoria; outputs/2026-06-11_resultados-eeg_v1.md.")
    slides_script.append(("13", "Responder a la profesora: si, integramos husos y memoria; no hacemos correlacion directa porque seria metodologicamente incorrecto con estas fuentes."))

    # 14
    s = new_slide(prs, 14, "ritmos", "Ritmos EEG esperados", "La interpretacion fisiologica se apoya en bandas, no solo en la etiqueta 'durmio'.")
    add_picture_fit(s, FIGURES / "2026-06-11_eeg_ritmos_v1.png", Inches(0.7), Inches(1.83), Inches(11.9), Inches(4.15))
    add_source(s, "Fuentes: clase 9 EEG; clase 3 sueño/memoria.")
    slides_script.append(("14", "Usar esta slide para mostrar manejo teorico de delta, theta, alfa, sigma y beta."))

    # 15
    s = new_slide(prs, 15, "moduladores", "Variables externas y no conciliacion", "Cronotipo, estres, cafeina y contexto pueden alterar la latencia y arquitectura del sueño.")
    add_table(
        s,
        [
            ["Variable", "Mecanismo visto en clase", "Aplicacion al caso"],
            ["Cronotipo / Proceso C", "fase circadiana modula propension al sueño", "siesta a horario no optimo puede fallar"],
            ["Proceso S / adenosina", "presion homeostatica aumenta con vigilia", "cafeina puede bloquear somnolencia"],
            ["Estres", "cortisol/NA modulan arousal y memoria", "recuperatorio o nervios pueden sostener alerta"],
            ["Ambiente/dispositivo", "higiene y comodidad afectan dormir", "laboratorio no habitual dificulta conciliacion"],
        ],
        Inches(0.65),
        Inches(1.8),
        Inches(12.0),
        Inches(3.6),
        [Inches(2.15), Inches(4.4), Inches(5.45)],
        font_size=10.5,
    )
    add_text(s, "Estas variables explican variacion y limitaciones; no son causas personales demostradas si no fueron medidas.", Inches(0.95), Inches(5.75), Inches(11.45), Inches(0.45), size=14.5, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_source(s, "Fuentes: clase 7 ritmos/higiene; clase 6 estres; clase 8 cafeina/farmacos.")
    slides_script.append(("15", "Conectar el no dormir con contenidos de clase, no con una explicacion anecdótica."))

    # 16
    s = new_slide(prs, 16, "discusion", "Discusion integrada", "Los resultados son coherentes con consolidacion activa, con limites metodologicos claros.")
    add_card(s, Inches(0.75), Inches(1.95), Inches(3.75), Inches(3.0), "Conducta", "El sujeto sueño rindio mejor en palabra objetivo; definicion fue alta en general.", BLUE)
    add_card(s, Inches(4.8), Inches(1.95), Inches(3.75), Inches(3.0), "Fisiologia", "El dataset secundario muestra NREM/N3 probable y ritmos analizables.", GREEN)
    add_card(s, Inches(8.85), Inches(1.95), Inches(3.75), Inches(3.0), "Teoria", "Ondas lentas, husos y ripples ofrecen mecanismo plausible de consolidacion.", ORANGE)
    add_text(s, "Interpretacion: H1 esta apoyada por el patron, pero la evidencia no es confirmatoria por tamaño muestral y fuentes separadas.", Inches(0.95), Inches(5.45), Inches(11.35), Inches(0.55), size=15, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_source(s, "Fuentes: outputs de memoria/EEG; clases 3, 7 y 8.")
    slides_script.append(("16", "Esta es la slide de discusion tipo paper: unir resultados con teoria y limites."))

    # 17
    s = new_slide(prs, 17, "limitaciones", "Limitaciones y proximos analisis", "El resultado permite discutir, no cerrar causalidad.")
    add_table(
        s,
        [
            ["Limitacion", "Impacto", "Como mejorarlo"],
            ["Sueño n=1", "no inferencia estadistica", "mas sujetos por condicion"],
            ["EEG secundario", "sin correlacion directa", "mismo sujeto: EEG + memoria"],
            ["Scoring tentativo", "clasificacion no clinica", "clave oficial / scoring experto"],
            ["Sin detector de husos", "sigma no equivale a husos", "detectar densidad y acoplamiento"],
        ],
        Inches(0.75),
        Inches(1.85),
        Inches(11.8),
        Inches(3.45),
        [Inches(3.0), Inches(4.1), Inches(4.7)],
        font_size=10.8,
    )
    add_text(s, "Proximo paso ideal: correlacionar densidad/acoplamiento de husos con mejora de memoria en el mismo sujeto.", Inches(0.95), Inches(5.75), Inches(11.4), Inches(0.45), size=14.5, bold=True, color=INK, align=PP_ALIGN.CENTER)
    add_source(s, "Fuentes: criterio metodologico del TP; clase 8 husos y memoria.")
    slides_script.append(("17", "Nombrar limitaciones como parte de la calidad academica, no como debilidad del grupo."))

    # 18
    s = new_slide(prs, 18, "conclusion", "Conclusion", "La hipotesis se apoya de forma descriptiva/parcial, no se confirma de forma causal fuerte.")
    add_bullets(
        s,
        [
            "Hipotesis: el sueño mejora la consolidacion de memoria declarativa.",
            "Resultado conductual: el sujeto en sueño tuvo mejor recuerdo formal de palabras.",
            "Mecanismo teorico: SWS/NREM, ondas lentas, husos y ripples pueden favorecer consolidacion.",
            "Decision: H1 queda apoyada descriptivamente; no se acepta como prueba estadistica concluyente.",
            "Variables externas ayudan a interpretar la no conciliacion y la variabilidad del sueño.",
        ],
        Inches(0.9),
        Inches(1.85),
        Inches(11.5),
        Inches(3.75),
        size=16.5,
    )
    add_card(s, Inches(1.15), Inches(5.85), Inches(10.9), Inches(0.62), "Cierre", "Formato paper: evidencia, mecanismo, limites y conclusion prudente.", ORANGE)
    add_source(s, "Fuentes: todos los outputs generados; consigna oficial; clases teoricas locales.")
    slides_script.append(("18", "Cerrar con una respuesta directa: se apoya parcialmente la hipotesis, con cautela metodologica."))

    prs.save(PPTX_OUT)
    return slides_script


def write_script(slides_script):
    lines = [
        "# Guion de presentacion v2 - sueño, memoria y consolidacion",
        "",
        "Fecha: 2026-06-11",
        "Duracion objetivo: 12-15 minutos",
        "Formato: paper oral en PPTX",
        "",
        "## Uso",
        "",
        "Este guion acompaña `outputs/2026-06-11_presentacion-sueno-memoria_v2.pptx`. La narrativa sigue estructura tipo paper: resumen, introduccion, hipotesis, metodos, resultados, discusion, limitaciones y conclusion.",
        "",
        "## Tesis oral",
        "",
        "La hipotesis principal es que el sueño mejora la consolidacion de memoria declarativa. Con los datos del practico, la hipotesis queda apoyada de forma descriptiva/parcial por el mejor rendimiento del sujeto en condicion sueño, pero no se confirma estadisticamente por el tamaño muestral y porque el EEG con sueño proviene de un dataset secundario.",
        "",
    ]
    for idx, text in slides_script:
        lines.extend([f"## Slide {idx}", "", text, ""])
    lines.extend(
        [
            "## Cierre metodologico obligatorio",
            "",
            "- No afirmar que el sujeto propio OpenBCI durmio.",
            "- No presentar el dataset secundario como si fuera del mismo sujeto/sesion.",
            "- Recordar que memoria en condicion sueño tiene n=1.",
            "- No calcular correlacion directa husos-memoria con fuentes separadas.",
            "- Decir que los husos/sigma se integran como mecanismo teorico y analisis fisiologico descriptivo.",
            "- Usar factores como cafeina, estres, pantallas, cronotipo y ambiente como variables plausibles, no causas demostradas.",
            "",
        ]
    )
    SCRIPT_OUT.write_text("\n".join(lines), encoding="utf-8")


def main() -> None:
    OUTPUTS.mkdir(exist_ok=True)
    slides_script = build_deck()
    write_script(slides_script)
    print(f"wrote {PPTX_OUT.relative_to(ROOT)}")
    print(f"wrote {SCRIPT_OUT.relative_to(ROOT)}")


if __name__ == "__main__":
    main()
