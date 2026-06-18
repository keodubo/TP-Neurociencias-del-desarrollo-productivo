#!/usr/bin/env python3
"""Arma presentacion.pdf y presentacion.pptx a partir de los PNG por slide
renderizados con Chrome headless (uno por diapositiva, 2560x1440).

Cada diapositiva se renderiza como slide "presente", de modo que el centrado
vertical y las ondas (canvas) salen exactamente como en el navegador.
"""
import glob
import os
import sys

import img2pdf
from pptx import Presentation
from pptx.util import Inches

SRC = sys.argv[1] if len(sys.argv) > 1 else "/tmp/deck_png"
DEST = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))  # presentacion-integrada-html/

pngs = sorted(glob.glob(os.path.join(SRC, "slide_*.png")))
if not pngs:
    sys.exit("No se encontraron PNGs en %s" % SRC)

# ---------- PDF (PNG embebido sin pérdida, página 16:9) ----------
pdf_path = os.path.join(DEST, "presentacion.pdf")
layout = img2pdf.get_layout_fun((img2pdf.in_to_pt(13.333), img2pdf.in_to_pt(7.5)))
with open(pdf_path, "wb") as f:
    f.write(img2pdf.convert(pngs, layout_fun=layout))
print("PDF  ->", pdf_path, "|", len(pngs), "páginas")

# ---------- PPTX (16:9, imagen a sangre completa por slide) ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]  # layout en blanco
for p in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(p, 0, 0, width=prs.slide_width, height=prs.slide_height)
pptx_path = os.path.join(DEST, "presentacion.pptx")
prs.save(pptx_path)
print("PPTX ->", pptx_path, "|", len(pngs), "slides")
